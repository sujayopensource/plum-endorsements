#!/usr/bin/env python3
"""
Generates Plum_Endorsement_Demo.pptx — a ~29-slide deck mirroring the video demo.

Usage:
    python3 video-demo/generate-ppt.py

Requirements:
    pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ────────────────────────────────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x0C, 0x29)
BG_MID = RGBColor(0x30, 0x2B, 0x63)
BG_LIGHT = RGBColor(0x24, 0x24, 0x3E)
ACCENT_1 = RGBColor(0x66, 0x7E, 0xEA)  # purple-blue
ACCENT_2 = RGBColor(0x76, 0x4B, 0xA2)  # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_70 = RGBColor(0xB3, 0xB3, 0xB3)
WHITE_50 = RGBColor(0x80, 0x80, 0x80)
WHITE_40 = RGBColor(0x66, 0x66, 0x66)
LABEL_COLOR = RGBColor(0xA5, 0xB4, 0xFC)  # light indigo
GREEN = RGBColor(0x34, 0xD3, 0x99)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
MONO_COLOR = RGBColor(0xC4, 0xB5, 0xFD)  # lavender for diagrams
CHECK_GREEN = RGBColor(0x4A, 0xDE, 0x80)

FONT_FAMILY = "Calibri"
MONO_FONT = "SF Mono"

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Plum_Endorsement_Demo.pptx")


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_rect(slide, left, top, width, height, color1=ACCENT_1, color2=ACCENT_2):
    """Add a rectangle with gradient fill as a decorative element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_FAMILY, line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    label_size=20, desc_size=18, spacing=14):
    """Add a bullet list with label: desc format.

    items: list of (label, description) tuples
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (label, desc) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Label run
        run_label = p.add_run()
        run_label.text = label + "  "
        run_label.font.size = Pt(label_size)
        run_label.font.color.rgb = LABEL_COLOR
        run_label.font.bold = True
        run_label.font.name = FONT_FAMILY

        # Description run
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(desc_size)
        run_desc.font.color.rgb = WHITE_70
        run_desc.font.bold = False
        run_desc.font.name = FONT_FAMILY

        p.space_after = Pt(spacing)
        p.space_before = Pt(2)

    return txBox


def add_mono_text(slide, left, top, width, height, text, font_size=13,
                  color=MONO_COLOR):
    """Add monospaced text for diagrams."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = MONO_FONT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = Pt(font_size * 1.45)
    return txBox


def add_section_number(slide, section_text):
    """Add a small section label in the top-left."""
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.4),
                section_text, font_size=12, color=WHITE_50)


def add_footer(slide, text="Plum Endorsement Management System"):
    """Add footer text at bottom."""
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.4),
                text, font_size=10, color=WHITE_40)


def make_title_slide(prs, title, subtitle, section_label=None):
    """Create a title card slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Accent bar
    add_gradient_rect(slide, Inches(4.5), Inches(2.6), Inches(4.3), Inches(0.06))

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.5),
                title, font_size=48, color=ACCENT_1, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(1.2),
                subtitle, font_size=24, color=WHITE_70,
                alignment=PP_ALIGN.CENTER)

    if section_label:
        add_section_number(slide, section_label)

    # Logo text
    add_textbox(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.5),
                "PLUM ENDORSEMENT MANAGEMENT SYSTEM",
                font_size=12, color=WHITE_40, alignment=PP_ALIGN.CENTER)

    return slide


def make_content_slide(prs, title, items, section_label=None):
    """Create a content slide with bullet items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title accent bar
    add_gradient_rect(slide, Inches(0.7), Inches(0.9), Inches(0.08), Inches(0.7))

    # Title
    add_textbox(slide, Inches(1.0), Inches(0.8), Inches(11), Inches(0.9),
                title, font_size=36, color=ACCENT_1, bold=True)

    # Bullet list
    add_bullet_list(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(5.0),
                    items, label_size=20, desc_size=17, spacing=12)

    if section_label:
        add_section_number(slide, section_label)
    add_footer(slide)

    return slide


# ── Slide Builders ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Large accent bar
    add_gradient_rect(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.08))

    # Main title
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                "Plum Endorsement\nManagement System",
                font_size=52, color=ACCENT_1, bold=True,
                alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(4.4), Inches(9.3), Inches(1.0),
                "Architecture  |  Live Demo  |  AI/Ollama  |  Vision",
                font_size=22, color=WHITE_70, alignment=PP_ALIGN.CENTER)

    # Bottom accent
    add_gradient_rect(slide, Inches(5.5), Inches(5.5), Inches(2.3), Inches(0.04))

    add_textbox(slide, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.5),
                "PLUM ENDORSEMENT MANAGEMENT SYSTEM",
                font_size=12, color=WHITE_40, alignment=PP_ALIGN.CENTER)


def slide_02_four_problems(prs):
    """Slide 2: The Four Hard Problems."""
    make_content_slide(prs, "The Four Hard Problems", [
        ("1. Coverage Gap",
         "Employee joins, endorsement created, but insurer takes days to confirm. "
         "During that window, the employee has no coverage."),
        ("2. Financial Drain",
         "Each ADD endorsement locks premium in an EA account. "
         "Without optimization, employers maintain a massive float."),
        ("3. Multi-Insurer Chaos",
         "4+ insurers, each with different APIs: REST, SOAP, CSV/SFTP. "
         "Different SLAs, different batch constraints."),
        ("4. Invisible Failures",
         "Submissions fail silently. Nobody knows until month-end reconciliation."),
    ], section_label="Section 1 - Problem Statement")


def slide_03_at_a_glance(prs):
    """Slide 3: System at a Glance."""
    make_content_slide(prs, "System at a Glance", [
        ("27 REST Endpoints", "5 controllers covering full endorsement lifecycle"),
        ("4 Insurer Integrations", "ICICI Lombard (REST), Niva Bupa (CSV/SFTP), Bajaj Allianz (SOAP), Mock"),
        ("13 Database Tables", "PostgreSQL 16 with Flyway migrations, optimistic locking"),
        ("4 Kafka Topics", "88 partitions, employer-ID partitioning, acks=all"),
        ("5 AI Modules", "2 with Ollama GenAI augmentation, 3 rule-based"),
        ("7 Grafana Dashboards", "40+ custom metrics, distributed tracing"),
        ("800+ Tests", "Unit + API + BDD + E2E + Performance. Zero failures."),
    ], section_label="Section 1 - Problem Statement")


def slide_04_c4_context(prs):
    """Slide 4: C4 Context Diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 2 - Architecture")

    add_gradient_rect(slide, Inches(0.7), Inches(0.7), Inches(0.08), Inches(0.7))
    add_textbox(slide, Inches(1.0), Inches(0.6), Inches(10), Inches(0.9),
                "C4 Model \u2014 System Context (Level 1)",
                font_size=32, color=ACCENT_1, bold=True)
    add_textbox(slide, Inches(1.0), Inches(1.3), Inches(10), Inches(0.5),
                "Who uses the system and what it connects to",
                font_size=16, color=WHITE_50)

    diagram = (
        "  \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510     "
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510     "
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "  \u2502  HR Admins    \u2502     \u2502  Finance      \u2502     \u2502  Ops Teams    \u2502\n"
        "  \u2502  Create &     \u2502     \u2502  EA Balance   \u2502     \u2502  Monitor &    \u2502\n"
        "  \u2502  manage       \u2502     \u2502  management   \u2502     \u2502  resolve      \u2502\n"
        "  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518     "
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518     "
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "         \u2502                   \u2502                    \u2502\n"
        "         \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
        "\u253c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "                             \u25bc\n"
        "              \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "              \u2502   PLUM ENDORSEMENT SYSTEM    \u2502\n"
        "              \u2502                               \u2502\n"
        "              \u2502  React UI + Spring Boot API   \u2502\n"
        "              \u2502  Kafka Events + PostgreSQL    \u2502\n"
        "              \u2502  5 Intelligence Pillars       \u2502\n"
        "              \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518\n"
        "                              \u2502\n"
        "         \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500"
        "\u253c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "         \u25bc                    \u25bc                    \u25bc\n"
        "  \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510  "
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510  "
        "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
        "  \u2502 ICICI Lombard \u2502  \u2502  Niva Bupa    \u2502  \u2502 Bajaj Allianz      \u2502\n"
        "  \u2502  REST/JSON    \u2502  \u2502  CSV/SFTP     \u2502  \u2502   SOAP/XML         \u2502\n"
        "  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518  "
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518  "
        "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518"
    )

    add_mono_text(slide, Inches(1.5), Inches(1.9), Inches(10.5), Inches(5.5),
                  diagram, font_size=12)
    add_footer(slide)


def slide_05_c4_container(prs):
    """Slide 5: C4 Container Diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 2 - Architecture")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "C4 Model \u2014 Container Diagram (Level 2)",
                font_size=30, color=ACCENT_1, bold=True)

    # Build two-column layout
    items_left = [
        ("React 19 SPA", "TanStack Table + shadcn/ui (:5173)"),
        ("Spring Boot 3.4", "Java 21 + Virtual Threads (:8080)"),
        ("5 Controllers", "27 REST Endpoints, RFC 7807"),
        ("3 CQRS Handlers", "8 Schedulers, 5 Services"),
        ("Ollama GenAI", "2 adapters deployed"),
    ]
    items_right = [
        ("PostgreSQL 16", "13 tables, Flyway migrations"),
        ("Redis 7", "Distributed cache, 60s TTL"),
        ("Kafka KRaft", "4 topics, 88 partitions"),
        ("Prometheus :9090", "15s scrape, 40+ metrics"),
        ("Grafana :3000", "7 dashboards, auto-provisioned"),
        ("Jaeger :16686", "100% sampling, trace propagation"),
        ("ELK Stack", "Elasticsearch + Logstash + Kibana"),
    ]

    add_textbox(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.5),
                "Application Tier", font_size=18, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.5),
                    items_left, label_size=17, desc_size=15, spacing=10)

    add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
                "Data & Observability", font_size=18, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(5.0),
                    items_right, label_size=17, desc_size=15, spacing=10)

    add_footer(slide)


def slide_06_hex_arch(prs):
    """Slide 6: Hexagonal Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 2 - Architecture")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "C4 Component \u2014 Hexagonal Architecture",
                font_size=30, color=ACCENT_1, bold=True)

    layers = [
        ("API Layer", "5 Controllers  \u00b7  27 Endpoints  \u00b7  RFC 7807 Error Handling",
         Inches(2.5), Inches(1.5), Inches(8.3), Inches(0.9), ACCENT_1),
        ("Application Layer",
         "3 Handlers (CQRS)  \u00b7  5 Services  \u00b7  8 Schedulers  |  Stateless  \u00b7  @Transactional",
         Inches(2.0), Inches(2.6), Inches(9.3), Inches(0.9), ACCENT_2),
        ("Domain Core",
         "Endorsement (11-state)  \u00b7  EAAccount  \u00b7  18 Ports  \u00b7  "
         "EndorsementEvent sealed (24 types)  |  >>> ZERO infrastructure imports <<<",
         Inches(1.5), Inches(3.7), Inches(10.3), Inches(1.1), GREEN),
        ("Infrastructure Layer",
         "JPA: 10 adapters + mappers   |   Insurer: Mock \u00b7 ICICI \u00b7 Niva \u00b7 Bajaj\n"
         "Kafka: 4 topics, 88 parts   |   Intelligence: 5 rule-based + 2 Ollama GenAI\n"
         "Resilience: Circuit Breakers + Retry",
         Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.8), GOLD),
    ]

    for name, desc, left, top, width, height, color in layers:
        # Background box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        # Layer name
        add_textbox(slide, left + Inches(0.2), top + Inches(0.05),
                    Inches(3), Inches(0.4),
                    name, font_size=16, color=color, bold=True)

        # Layer description
        add_textbox(slide, left + Inches(0.2), top + Inches(0.35),
                    width - Inches(0.4), height - Inches(0.4),
                    desc, font_size=13, color=WHITE_70)

    add_footer(slide)


def slide_07_patterns_tech(prs):
    """Slide 7: Design Patterns & Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 2 - Architecture")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Design Patterns & Technology Stack",
                font_size=30, color=ACCENT_1, bold=True)

    patterns = [
        ("Strategy", "InsurerPort \u2014 add insurer = add class + DB row"),
        ("State", "11-state lifecycle with canTransitionTo()"),
        ("Observer", "EventPublisher \u2192 Kafka \u2192 consumers"),
        ("Factory", "InsurerRouter.resolve() from DB config"),
        ("Adapter", "Domain never imports infrastructure"),
        ("CQRS", "Command handlers vs read-only QueryHandler"),
    ]

    tech = [
        ("Java 21", "Virtual Threads \u2014 1M+ concurrent"),
        ("Spring Boot 3.4", "Actuator, virtual thread support"),
        ("PostgreSQL 16", "ACID, optimistic locking"),
        ("Kafka KRaft", "32 partitions, employerId key"),
        ("Resilience4j", "Per-insurer circuit breakers"),
        ("ZGC", "Sub-millisecond GC pauses"),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
                "Design Patterns", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
                    patterns, label_size=17, desc_size=15, spacing=8)

    add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
                "Technology Stack", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0),
                    tech, label_size=17, desc_size=15, spacing=8)

    add_footer(slide)


def slide_08_lifecycle(prs):
    """Slide 8: Complete Endorsement Lifecycle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 3 - Lifecycle Demo")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Complete Endorsement Lifecycle",
                font_size=30, color=ACCENT_1, bold=True)

    # Flow diagram
    states = [
        "CREATED", "VALIDATED", "PROVISIONALLY\nCOVERED",
        "SUBMITTED\nTO INSURER", "INSURER\nPROCESSING", "CONFIRMED"
    ]
    colors = [WHITE_70, WHITE_70, GREEN, ACCENT_1, ACCENT_2, CHECK_GREEN]

    start_x = Inches(0.5)
    y = Inches(2.3)
    box_w = Inches(1.8)
    box_h = Inches(1.0)
    gap = Inches(0.3)

    for i, (state, color) in enumerate(zip(states, colors)):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = state
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

        # Arrow between boxes
        if i < len(states) - 1:
            arrow_x = x + box_w
            add_textbox(slide, arrow_x, y + Inches(0.3), gap, Inches(0.4),
                        "\u2192", font_size=20, color=WHITE_50,
                        alignment=PP_ALIGN.CENTER)

    # Key points below
    items = [
        ("Idempotency", "Duplicate check via unique key at CREATED"),
        ("Provisional Coverage", "Employee covered immediately at step 3, not at confirmation"),
        ("Strategy Pattern", "InsurerRouter resolves adapter; handler never knows the protocol"),
        ("State Machine", "canTransitionTo() validates every transition at compile time"),
        ("Batch Path", "For batch-only insurers: Queued \u2192 Batch Submitted \u2192 Confirmed"),
    ]
    add_bullet_list(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(3.5),
                    items, label_size=17, desc_size=15, spacing=8)

    add_footer(slide)


def slide_09_provisional(prs):
    """Slide 9: Provisional Coverage."""
    make_content_slide(prs, "Provisional Coverage: Key Innovation", [
        ("The Problem",
         "Insurer confirmation takes 1\u20135 days. Employee has no coverage during that window."),
        ("Our Solution",
         "Provisional coverage activated at endorsement creation, before insurer response."),
        ("30-Day Safety Net",
         "Provisional coverage auto-expires after 30 days if insurer never confirms."),
        ("Upgrade on Confirmation",
         "When insurer confirms, provisional \u2192 confirmed. Seamless for employee."),
        ("Batch Path",
         "For batch-only insurers (Niva Bupa): Queued for Batch \u2192 Batch Submitted every "
         "15 min \u2192 Confirmed. Coverage active the entire time."),
    ], section_label="Section 3 - Lifecycle Demo")


def slide_10_ea_balance(prs):
    """Slide 10: EA Balance Optimization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 4 - EA Balance")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "EA Balance Optimization",
                font_size=30, color=ACCENT_1, bold=True)

    # Math section
    math_items = [
        ("Naive Approach", "60 ADDs \u00d7 \u20b91,000 = \u20b960,000 peak balance required"),
        ("With Optimization", "Process 40 DELETEs first \u2192 +\u20b932,000 credit released"),
        ("Net Requirement", "\u20b928,000 \u2014 53% savings in float"),
        ("Algorithm", "0/1 Knapsack with dynamic programming, 15-minute scheduler"),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
                "Optimization Math", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0),
                    math_items, label_size=17, desc_size=15, spacing=10)

    # Scoring section
    scoring_items = [
        ("Priority Order", "DELETEs first (free balance) \u2192 cost-neutral UPDATEs \u2192 ADDs by urgency"),
        ("Composite Score", "Urgency weight 60% + EA impact weight 40%"),
        ("EA Lookup", "3 KPI cards: Total Balance, Reserved Amount, Available Balance"),
    ]

    add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
                "Scoring & Visibility", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0),
                    scoring_items, label_size=17, desc_size=15, spacing=10)

    # Big number callout
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.0), Inches(5.0), Inches(5.3), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    add_textbox(slide, Inches(4.3), Inches(5.1), Inches(4.7), Inches(0.6),
                "53% Reduction in EA Float",
                font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.3), Inches(5.7), Inches(4.7), Inches(0.6),
                "\u20b960,000 \u2192 \u20b928,000 through priority ordering",
                font_size=16, color=WHITE_70, alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_11_visibility(prs):
    """Slide 11: Real-Time Visibility Screens."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 5 - Visibility")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Real-Time Visibility: 10 UI Screens",
                font_size=30, color=ACCENT_1, bold=True)

    screens_left = [
        ("Dashboard", "KPI cards: total, pending, confirmed, failed + EA balance"),
        ("Create Endorsement", "Progressive disclosure form, adapts by type"),
        ("Endorsement Detail", "Status timeline, insurer reference, full history"),
        ("Endorsement List", "Sortable TanStack Table, status filters, CSV export"),
        ("Batch Progress", "Batch status: Submitted, Completed, Failed"),
    ]

    screens_right = [
        ("Insurer Config", "4 insurers with capabilities, rate limits"),
        ("EA Account Lookup", "Balance, reserved, available with progress bar"),
        ("Reconciliation", "Matched, Partial, Rejected, Missing counts"),
        ("Intelligence Hub", "5 tabs: Anomalies, Forecasts, Errors, Process, Health"),
        ("Audit Log", "Full event history with filters"),
    ]

    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                    screens_left, label_size=17, desc_size=14, spacing=8)
    add_bullet_list(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(5.5),
                    screens_right, label_size=17, desc_size=14, spacing=8)

    add_footer(slide)


def slide_12_intelligence_overview(prs):
    """Slide 12: Intelligence & AI Overview."""
    make_content_slide(prs, "Intelligence & AI: 5 Pillars", [
        ("Anomaly Detection",
         "5 rules (Volume Spike, ADD/DELETE Cycling, Suspicious Timing, "
         "Unusual Premium, Dormancy Break)"),
        ("Balance Forecasting",
         "30-day EA projection with dual seasonality (day-of-week + monthly)"),
        ("Error Resolution",
         "5 error patterns with confidence scoring, auto-apply above 95%"),
        ("Process Mining",
         "STP rate tracking, per-insurer optimization, transition metrics"),
        ("Employer Health Score",
         "Composite score from endorsement patterns, EA utilization, error rates"),
    ], section_label="Section 6 - Intelligence")


def slide_13_ollama(prs):
    """Slide 13: How We Use Ollama."""
    make_content_slide(prs, "How We Use Ollama", [
        ("Runtime",
         "Docker container running llama3.2 model locally \u2014 zero cloud API costs"),
        ("Integration",
         "Spring AI Ollama starter with ChatClient.Builder"),
        ("Activation",
         "@ConditionalOnProperty \u2014 activated via 'ollama' Spring profile"),
        ("Deployed Adapters",
         "OllamaAugmentedAnomalyDetector + OllamaErrorResolver"),
        ("Resilience",
         "@CircuitBreaker + @Retry with graceful fallback to rule-based results"),
        ("Config",
         "Temperature 0.3 for deterministic output, 512 token max"),
    ], section_label="Section 6 - Intelligence")


def slide_14_anomaly(prs):
    """Slide 14: Anomaly Detection."""
    make_content_slide(prs, "Anomaly Detection", [
        ("Volume Spike",
         "Detects unusual endorsement volume for an employer"),
        ("ADD/DELETE Cycling",
         "Flags repeated add-then-delete patterns suggesting premium arbitrage"),
        ("Suspicious Timing",
         "Off-hours or weekend submissions flagged"),
        ("Unusual Premium",
         "Premium amounts outside historical norms"),
        ("Dormancy Break",
         "Sudden activity from long-dormant employers"),
        ("Ollama Enrichment",
         "Anomalies scoring > 0.7 sent to LLM for human-readable explanation "
         "and action recommendation. Narrative layer \u2014 enriches, not decides."),
    ], section_label="Section 6 - Intelligence")


def slide_15_forecast_error(prs):
    """Slide 15: Forecasting & Error Resolution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 6 - Intelligence")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Forecasting & Error Resolution",
                font_size=30, color=ACCENT_1, bold=True)

    forecast_items = [
        ("30-Day Projection", "EA balance forecast using historical data"),
        ("Dual Seasonality", "Day-of-week (Mon=1.2x) + monthly (Apr=1.4x fiscal year)"),
        ("Top-Up Alerts", "Shortfall detection triggers proactive notifications"),
    ]

    error_items = [
        ("5 Error Patterns", "Date format, missing field, duplicate, invalid amount, timeout"),
        ("Confidence Scoring", "Above 95% \u2192 auto-applied; below \u2192 human approval"),
        ("Success Tracking", "Feedback loop tracks auto-resolution success/failure rates"),
        ("Ollama Resolver",
         "Ambiguous cases (<95%) analyzed by LLM for specific fix with reasoning"),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
                "Balance Forecasting", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0),
                    forecast_items, label_size=17, desc_size=15, spacing=10)

    add_textbox(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
                "Error Resolution", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0),
                    error_items, label_size=17, desc_size=15, spacing=10)

    add_footer(slide)


def slide_16_process_mining(prs):
    """Slide 16: Process Mining & 3-Stage Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 6 - Intelligence")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Process Mining & 3-Stage AI Evolution",
                font_size=30, color=ACCENT_1, bold=True)

    mining_items = [
        ("STP Rate", "% of endorsements completing without human intervention"),
        ("Per-Insurer Cards", "Focus optimization on lowest-performing insurers"),
        ("30-Day Trend", "STP trajectory chart showing improvement over time"),
        ("Transition Metrics", "Measures every state change for bottleneck detection"),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
                "Process Mining", font_size=20, color=GOLD, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11.3), Inches(2.5),
                    mining_items, label_size=17, desc_size=15, spacing=8)

    # Stage comparison table
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
                "3-Stage Evolution Comparison", font_size=20, color=GOLD, bold=True)

    stages = [
        ("Stage 1: Rule-Based", "5 adapters deployed, zero ML deps, production-ready",
         WHITE_70),
        ("Stage 2: Ollama GenAI",
         "2 deployed (anomaly + error), 3 planned. Local LLM, no cloud costs",
         GREEN),
        ("Stage 3: Full ML",
         "Isolation Forest, Prophet, RAG, PM4Py, OR-Tools. Adapter swap behind same port",
         ACCENT_1),
    ]

    for i, (stage, desc, color) in enumerate(stages):
        y = Inches(4.6) + i * Inches(0.85)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(1.0), y, Inches(11.0), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.35),
                    stage, font_size=15, color=color, bold=True)
        add_textbox(slide, Inches(1.2), y + Inches(0.35), Inches(10.5), Inches(0.35),
                    desc, font_size=13, color=WHITE_70)

    add_footer(slide)


def slide_17_docker_services(prs):
    """Slide 17: 9 Docker Services."""
    make_content_slide(prs, "Observability: 9 Docker Services", [
        ("PostgreSQL :5432", "13 tables, Flyway migrations, ACID"),
        ("Redis :6379", "Distributed cache, 60s TTL, @Cacheable"),
        ("Kafka :9092", "KRaft mode, 4 topics, 88 partitions"),
        ("Prometheus :9090", "15s scrape interval, 40+ custom metrics"),
        ("Grafana :3000", "7 auto-provisioned dashboards"),
        ("Jaeger :16686", "100% sampling, distributed tracing"),
        ("Elasticsearch :9200", "Log aggregation, structured JSON"),
        ("Logstash :5000", "Log pipeline with TCP input"),
        ("Kibana :5601", "Log visualization and search"),
    ], section_label="Section 7 - Observability")


def slide_18_grafana(prs):
    """Slide 18: Grafana Dashboards."""
    make_content_slide(prs, "Grafana Dashboards", [
        ("Application Overview",
         "Request rate, P95 latency, error rate, JVM heap, threads, DB pool"),
        ("Business Metrics",
         "Creation rate by type, endorsements by status, insurer submission latency "
         "(Bajaj 250ms vs ICICI 150ms)"),
        ("Intelligence Monitoring",
         "Anomaly detection rate, forecast shortfall detections, "
         "error auto-resolution gauge, STP trend, batch optimization savings"),
        ("Infrastructure Health",
         "Kafka consumer lag, PostgreSQL connections, Redis hit rate"),
        ("Multi-Insurer",
         "Per-insurer submission rates, circuit breaker states, SLA compliance"),
        ("Reconciliation",
         "Match rates, discrepancy trends, resolution time"),
        ("Scheduler",
         "Batch assembly timing, anomaly scan duration, forecast generation"),
    ], section_label="Section 7 - Observability")


def slide_19_tracing(prs):
    """Slide 19: Distributed Tracing."""
    make_content_slide(prs, "Distributed Tracing: Jaeger", [
        ("100% Sampling", "Every request traced end-to-end, no sampling loss"),
        ("OpenTelemetry", "Micrometer tracing bridge with OTLP exporter"),
        ("MDC Propagation",
         "traceId, spanId, requestId, endorsementId, employerId in every log line"),
        ("Correlation",
         "Search by endorsement ID to see full request journey across services"),
        ("Structured Logging",
         "JSON format in production via Logstash encoder, "
         "human-readable in development"),
    ], section_label="Section 7 - Observability")


def slide_20_scalability(prs):
    """Slide 20: Scalability."""
    make_content_slide(prs, "Scalability: 1M Endorsements/Day", [
        ("Stateless Services",
         "All private final fields. No mutable instance state. Horizontal scaling."),
        ("Java 21 Virtual Threads",
         "1M+ concurrent operations, no thread pool tuning required"),
        ("Kafka Partitioning",
         "32 partitions per topic, employer-ID key for per-employer ordering"),
        ("Optimistic Locking",
         "Version-based concurrency on Endorsement and EAAccount entities"),
        ("Idempotency Keys",
         "UNIQUE constraint, safe retries across network failures"),
        ("K8s HPA",
         "Auto-scale from 2 to 8 pods at 70% CPU threshold"),
    ], section_label="Section 8 - Scalability")


def slide_21_circuit_breakers(prs):
    """Slide 21: Per-Insurer Circuit Breakers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 8 - Scalability")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Per-Insurer Circuit Breakers",
                font_size=30, color=ACCENT_1, bold=True)

    # Table-style comparison
    headers = [
        ("Insurer", Inches(2.5)),
        ("Failure Threshold", Inches(2.2)),
        ("Window Size", Inches(2.0)),
        ("Wait (Open)", Inches(2.0)),
        ("Retry", Inches(2.0)),
    ]

    rows = [
        ("ICICI Lombard", "50%", "20 calls", "30s", "3 attempts, exp backoff"),
        ("Bajaj Allianz", "40%", "15 calls", "45s", "5 attempts, exp backoff"),
        ("Mock (default)", "50%", "10 calls", "30s", "3 attempts, exp backoff"),
    ]

    table_left = Inches(1.0)
    table_top = Inches(1.8)
    row_height = Inches(0.7)
    col_widths = [h[1] for h in headers]

    # Header row
    x = table_left
    for (name, w) in headers:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, w, row_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x20, 0x1E, 0x50)
        shape.line.color.rgb = ACCENT_2
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = GOLD
        p.font.bold = True
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        x += w

    # Data rows
    for ri, row in enumerate(rows):
        y = table_top + (ri + 1) * row_height
        x = table_left
        for ci, (val, (_, w)) in enumerate(zip(row, headers)):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x38)
            shape.line.color.rgb = RGBColor(0x30, 0x2E, 0x55)
            shape.line.width = Pt(0.5)

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE_70 if ci > 0 else LABEL_COLOR
            p.font.bold = ci == 0
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            x += w

    # Key points
    items = [
        ("Typed Fallback",
         "Every @CircuitBreaker method has a typed fallback that returns same type + "
         "Throwable. Never throws from a fallback."),
        ("Per-Insurer Tuning",
         "Each insurer tuned to its reliability profile. SOAP (Bajaj) gets lower "
         "threshold and longer wait."),
        ("Actuator Visibility",
         "/actuator/circuitbreakers exposes real-time state. Alert on OPEN."),
    ]
    add_bullet_list(slide, Inches(1.0), Inches(4.7), Inches(10.7), Inches(2.5),
                    items, label_size=17, desc_size=15, spacing=8)

    add_footer(slide)


def slide_22_tests(prs):
    """Slide 22: Test Coverage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 9 - Tests")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Test Coverage: 800+ Tests",
                font_size=30, color=ACCENT_1, bold=True)

    # Test pyramid
    test_levels = [
        ("Unit Tests", "420", "JUnit 5 + Mockito + AssertJ",
         Inches(3.5), Inches(6.3)),
        ("API Tests", "124", "REST Assured + Testcontainers",
         Inches(3.0), Inches(7.3)),
        ("BDD Tests", "92", "Cucumber \u2014 22 feature files",
         Inches(2.5), Inches(8.3)),
        ("E2E Tests", "158", "Playwright \u2014 30 spec files",
         Inches(2.0), Inches(9.3)),
        ("Performance", "6", "Gatling simulations",
         Inches(1.5), Inches(10.3)),
    ]

    for i, (name, count, tech, left, width) in enumerate(test_levels):
        y = Inches(1.5) + i * Inches(1.05)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, y, width, Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = ACCENT_1 if i == 0 else ACCENT_2
        shape.line.width = Pt(1.5)

        # Count
        add_textbox(slide, left + Inches(0.2), y + Inches(0.05),
                    Inches(1.2), Inches(0.4),
                    count, font_size=22, color=GREEN, bold=True)

        # Name
        add_textbox(slide, left + Inches(1.4), y + Inches(0.05),
                    Inches(2.5), Inches(0.4),
                    name, font_size=16, color=LABEL_COLOR, bold=True)

        # Tech
        add_textbox(slide, left + Inches(1.4), y + Inches(0.4),
                    width - Inches(1.8), Inches(0.4),
                    tech, font_size=13, color=WHITE_70)

    # Total
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.0), Inches(6.8), Inches(5.3), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
    shape.line.color.rgb = CHECK_GREEN
    shape.line.width = Pt(2)

    add_textbox(slide, Inches(4.3), Inches(6.85), Inches(4.7), Inches(0.4),
                "TOTAL: 800+   |   ZERO FAILURES",
                font_size=18, color=CHECK_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_23_allure(prs):
    """Slide 23: Allure Combined Report."""
    make_content_slide(prs, "Allure Combined Report", [
        ("Docker Service",
         "Allure Docker Service at :5050 with combined report"),
        ("API Section",
         "124 REST Assured tests with request/response payloads"),
        ("BDD Section",
         "92 Cucumber scenarios with step-by-step execution"),
        ("E2E Section",
         "158 Playwright tests with screenshots on failure"),
        ("Performance Section",
         "6 Gatling simulations with response time distributions"),
        ("Single Command",
         "./run-all-tests.sh builds, runs all suites, and publishes combined report"),
    ], section_label="Section 9 - Tests")


def slide_24_ai_vision(prs):
    """Slide 24: AI Automation Vision."""
    make_content_slide(prs, "AI Automation Vision", [
        ("Anomaly \u2192 ML",
         "Isolation Forest unsupervised learning \u2192 Autoencoder 2nd-stage "
         "\u2192 Graph-based fraud network detection"),
        ("Forecasting \u2192 ML",
         "Facebook Prophet multi-seasonality \u2192 LSTM neural networks "
         "\u2192 Ensemble methods"),
        ("Error Resolution \u2192 RAG",
         "Retrieval-Augmented Generation with vector database of past resolutions"),
        ("Process Mining \u2192 PM4Py",
         "Conformance checking via Python sidecar service"),
        ("Batch Optimization \u2192 OR-Tools",
         "Google OR-Tools linear programming replacing knapsack heuristic"),
    ], section_label="Section 10 - Vision")


def slide_25_ml_rollout(prs):
    """Slide 25: ML Rollout Pattern."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 10 - Vision")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "ML Rollout Pattern",
                font_size=30, color=ACCENT_1, bold=True)

    phases = [
        ("Phase 1: Implement", "New ML adapter behind existing port interface", ACCENT_1),
        ("Phase 2: Deploy", "Feature flag via @ConditionalOnProperty", ACCENT_2),
        ("Phase 3: Shadow", "Run alongside current adapter, compare outputs", GOLD),
        ("Phase 4: Promote", "Switch when metrics meet targets", GREEN),
    ]

    for i, (phase, desc, color) in enumerate(phases):
        y = Inches(1.5) + i * Inches(1.2)
        x = Inches(1.5) + i * Inches(0.3)
        w = Inches(10.0) - i * Inches(0.6)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, w, Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        add_textbox(slide, x + Inches(0.3), y + Inches(0.08),
                    Inches(4), Inches(0.4),
                    phase, font_size=18, color=color, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.45),
                    w - Inches(0.6), Inches(0.4),
                    desc, font_size=15, color=WHITE_70)

    # Cross-cutting
    items = [
        ("Cross-Cutting", "Feature Store, MLOps pipeline, A/B testing framework, "
         "natural language query interface"),
        ("Already Proven", "Pattern validated with 2 deployed Ollama adapters. "
         "Each stage is an adapter swap \u2014 hexagonal architecture payoff."),
    ]
    add_bullet_list(slide, Inches(1.0), Inches(6.0), Inches(11.0), Inches(1.5),
                    items, label_size=17, desc_size=15, spacing=8)

    add_footer(slide)


def slide_26_product_vision(prs):
    """Slide 26: Product Evolution Vision."""
    make_content_slide(prs, "Product Evolution Vision: Phase 4 Global Platform", [
        ("Advanced Analytics",
         "Complete intelligence layer to 100% coverage"),
        ("Multi-Currency EA",
         "International employer support with currency conversion"),
        ("Localized Integrations",
         "New market insurer integrations with country-specific protocols"),
        ("Regulatory Compliance",
         "Country-specific compliance rules as pluggable adapters"),
        ("Self-Service Portal",
         "Insurer onboarding without engineering involvement"),
        ("API Marketplace",
         "Platform API for third-party integrations"),
        ("Multi-Region Deployment",
         "Global low-latency deployment across regions"),
        ("Auth via Hexagonal",
         "JWT + OAuth2 with RBAC, designed as port + adapter pair"),
    ], section_label="Section 10 - Vision")


def slide_27_arch_roadmap(prs):
    """Slide 27: 5-Year Architecture Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_section_number(slide, "Section 10 - Vision")

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Architectural Vision: 5-Year Roadmap",
                font_size=30, color=ACCENT_1, bold=True)

    years = [
        ("2026", "Modular Monolith Hardening",
         "Fitness functions, distributed locking (ShedLock), "
         "automated architecture checks in CI", ACCENT_1),
        ("2027", "Service Extraction",
         "Strangler Fig pattern, multi-region CQRS read replicas, "
         "event-driven decomposition", ACCENT_2),
        ("2028", "Platform Maturity",
         "API marketplace, self-service insurer onboarding, "
         "partner ecosystem", GOLD),
        ("2029", "Intelligence Platform",
         "Data Mesh architecture, ML pipeline with feature store, "
         "real-time analytics", GREEN),
        ("2030", "Global Scale",
         "Active-active multi-region, autonomous operations, "
         "self-healing infrastructure", CHECK_GREEN),
    ]

    for i, (year, title, desc, color) in enumerate(years):
        y = Inches(1.4) + i * Inches(1.15)

        # Year badge
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(1.0), y, Inches(1.3), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x40)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        add_textbox(slide, Inches(1.1), y + Inches(0.2), Inches(1.1), Inches(0.5),
                    year, font_size=22, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Content
        add_textbox(slide, Inches(2.6), y + Inches(0.05), Inches(9), Inches(0.4),
                    title, font_size=17, color=color, bold=True)
        add_textbox(slide, Inches(2.6), y + Inches(0.4), Inches(9), Inches(0.5),
                    desc, font_size=14, color=WHITE_70)

        # Connecting line
        if i < len(years) - 1:
            line_y = y + Inches(0.9)
            add_textbox(slide, Inches(1.55), line_y, Inches(0.3), Inches(0.3),
                        "\u2502", font_size=16, color=WHITE_40,
                        alignment=PP_ALIGN.CENTER)

    # Thesis
    add_textbox(slide, Inches(1.0), Inches(7.0), Inches(11), Inches(0.4),
                "Thesis: Evolve from modular monolith to selectively-decomposed "
                "service platform, guided by automated fitness functions.",
                font_size=12, color=WHITE_50)


def slide_28_deliverables(prs):
    """Slide 28: Deliverable Mapping."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_gradient_rect(slide, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.6))
    add_textbox(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.8),
                "Deliverable Mapping",
                font_size=30, color=ACCENT_1, bold=True)

    deliverables = [
        ("1. Architecture",
         "Hexagonal architecture with C4 model views at Context, Container, "
         "and Component levels"),
        ("2. No Coverage Gap",
         "Provisional coverage at creation with 30-day safety net. "
         "Employee covered immediately."),
        ("3. EA Minimization",
         "Priority ordering with 53% savings (\u20b960K \u2192 \u20b928K). "
         "0/1 Knapsack DP algorithm."),
        ("4. Real-Time Visibility",
         "10 UI screens + 7 Grafana dashboards + distributed tracing"),
        ("5. AI/Automation",
         "5 intelligence pillars, 2 Ollama GenAI adapters deployed, "
         "3 planned. Full ML vision documented."),
        ("6. Working Prototype",
         "Live demo + 800+ tests + comprehensive documentation. "
         "Single-command startup."),
    ]

    for i, (title, desc) in enumerate(deliverables):
        y = Inches(1.4) + i * Inches(0.95)

        # Checkmark
        add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
                    "\u2713", font_size=22, color=CHECK_GREEN, bold=True)

        # Title
        add_textbox(slide, Inches(1.3), y + Inches(0.0), Inches(3), Inches(0.4),
                    title, font_size=17, color=LABEL_COLOR, bold=True)

        # Description
        add_textbox(slide, Inches(1.3), y + Inches(0.38), Inches(10.5), Inches(0.5),
                    desc, font_size=14, color=WHITE_70)

    add_footer(slide)


def slide_29_thank_you(prs):
    """Slide 29: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_gradient_rect(slide, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.06))

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
                "Thank You",
                font_size=56, color=ACCENT_1, bold=True,
                alignment=PP_ALIGN.CENTER)

    highlights = [
        "Production-Grade Code",
        "Hexagonal Architecture",
        "Cloud-Native Patterns",
        "Ollama GenAI Augmentation",
        "800+ Tests, Zero Failures",
        "5-Year Architectural Vision",
    ]

    add_textbox(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(0.5),
                "  \u00b7  ".join(highlights),
                font_size=16, color=WHITE_70, alignment=PP_ALIGN.CENTER)

    add_gradient_rect(slide, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.04))

    add_textbox(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.5),
                "PLUM ENDORSEMENT MANAGEMENT SYSTEM",
                font_size=12, color=WHITE_40, alignment=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating slides...")

    slide_01_title(prs)
    print("  [1/29] Title slide")

    slide_02_four_problems(prs)
    print("  [2/29] Four Hard Problems")

    slide_03_at_a_glance(prs)
    print("  [3/29] System at a Glance")

    slide_04_c4_context(prs)
    print("  [4/29] C4 Context")

    slide_05_c4_container(prs)
    print("  [5/29] C4 Container")

    slide_06_hex_arch(prs)
    print("  [6/29] Hexagonal Architecture")

    slide_07_patterns_tech(prs)
    print("  [7/29] Design Patterns & Tech")

    slide_08_lifecycle(prs)
    print("  [8/29] Endorsement Lifecycle")

    slide_09_provisional(prs)
    print("  [9/29] Provisional Coverage")

    slide_10_ea_balance(prs)
    print("  [10/29] EA Balance Optimization")

    slide_11_visibility(prs)
    print("  [11/29] Real-Time Visibility")

    slide_12_intelligence_overview(prs)
    print("  [12/29] Intelligence Overview")

    slide_13_ollama(prs)
    print("  [13/29] How We Use Ollama")

    slide_14_anomaly(prs)
    print("  [14/29] Anomaly Detection")

    slide_15_forecast_error(prs)
    print("  [15/29] Forecasting & Error Resolution")

    slide_16_process_mining(prs)
    print("  [16/29] Process Mining & 3-Stage")

    slide_17_docker_services(prs)
    print("  [17/29] Docker Services")

    slide_18_grafana(prs)
    print("  [18/29] Grafana Dashboards")

    slide_19_tracing(prs)
    print("  [19/29] Distributed Tracing")

    slide_20_scalability(prs)
    print("  [20/29] Scalability")

    slide_21_circuit_breakers(prs)
    print("  [21/29] Circuit Breakers")

    slide_22_tests(prs)
    print("  [22/29] Test Coverage")

    slide_23_allure(prs)
    print("  [23/29] Allure Report")

    slide_24_ai_vision(prs)
    print("  [24/29] AI Automation Vision")

    slide_25_ml_rollout(prs)
    print("  [25/29] ML Rollout Pattern")

    slide_26_product_vision(prs)
    print("  [26/29] Product Evolution")

    slide_27_arch_roadmap(prs)
    print("  [27/29] Architecture Roadmap")

    slide_28_deliverables(prs)
    print("  [28/29] Deliverable Mapping")

    slide_29_thank_you(prs)
    print("  [29/29] Thank You")

    prs.save(OUTPUT_FILE)
    size_kb = os.path.getsize(OUTPUT_FILE) / 1024
    print(f"\nSaved: {OUTPUT_FILE}")
    print(f"Size: {size_kb:.0f} KB")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
